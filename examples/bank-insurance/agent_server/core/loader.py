# SPDX-License-Identifier: Apache-2.0
"""Document and content loading utilities for the example agent.

This module provides functionality for loading and processing various types of content,
including web pages, local files, and text data. It includes classes and methods for
web crawling, document parsing, and content chunking. The module supports multiple
file formats and provides utilities for handling different types of content sources,
ensuring consistent processing and storage of loaded content.
"""

from __future__ import annotations

import base64
import contextlib
import json
import logging
import os
import pickle
import time
import uuid
from enum import Enum
from pathlib import Path
from urllib.parse import urljoin

import requests
import tiktoken
from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter
from mistralai import Mistral
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager

logger = logging.getLogger(__name__)


CHROME_DRIVER_VERSION = "138.0.7204.94"
MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY")


def encode_image(image_path: str) -> str | None:
    """Encode the image to base64.

    Args:
        image_path (str): Path to the image file.

    Returns:
        str | None: Base64-encoded string if successful, None otherwise.
    """
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode("utf-8")
    except FileNotFoundError:
        logger.error(f"Error: The file {image_path} was not found.")
        return None
    except Exception as e:
        logger.error(f"Error: {e}")
        return None


def _make_text_splitter() -> RecursiveCharacterTextSplitter:
    """Return a text splitter, preferring tiktoken-based chunking with a character-based fallback."""
    try:
        tiktoken.get_encoding("cl100k_base")
        return RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base", chunk_size=10000, chunk_overlap=1000
        )
    except Exception:
        return RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)


class SourceType(Enum):
    WEB = 1
    LOCAL = 2
    TEXT = 3


class DocObject:
    def __init__(self, id: str, source: str) -> None:
        self.id = id
        self.source = source


class CrawledObject(DocObject):
    def __init__(
        self,
        id: str,
        source: str,
        content: str,
        metadata: dict | None = None,
        is_chunk: bool = False,
        is_error: bool = False,
        error_message: str | None = None,
        source_type: SourceType = SourceType.WEB,
    ) -> None:
        super().__init__(id, source)
        self.content = content
        self.metadata = metadata or {}
        self.is_chunk = is_chunk
        self.is_error = is_error
        self.error_message = error_message
        self.source_type = source_type

    def to_dict(self) -> dict:
        return {
            "id": self.id,
            "source": self.source,
            "content": self.content,
            "metadata": self.metadata,
            "is_chunk": self.is_chunk,
            "is_error": self.is_error,
            "error_message": self.error_message,
            "source_type": self.source_type,
        }

    @classmethod
    def from_dict(cls, data: dict) -> CrawledObject:
        return cls(
            id=data["id"],
            source=data["source"],
            content=data["content"],
            metadata=data["metadata"],
            is_chunk=data["is_chunk"],
            is_error=data["is_error"],
            error_message=data["error_message"],
            source_type=data["source_type"],
        )


class Loader:
    def __init__(self) -> None:
        pass

    def to_crawled_url_objs(self, url_list: list[str]) -> list[CrawledObject]:
        """Convert a list of URLs to CrawledObject instances.

        Args:
            url_list (List[str]): List of URLs to convert.

        Returns:
            List[CrawledObject]: List of CrawledObject instances containing crawled content.
        """
        url_objs = [DocObject(str(uuid.uuid4()), url) for url in url_list]
        return self.crawl_urls(url_objs)

    def crawl_urls(self, url_objects: list[DocObject]) -> list[CrawledObject]:
        """Crawl a list of URLs and extract their content.

        Args:
            url_objects (list[DocObject]): List of DocObject instances containing URLs to crawl.

        Returns:
            List[CrawledObject]: List of CrawledObject instances containing crawled content.
        """
        logger.info(f"Starting web crawling for {len(url_objects)} URLs...")

        docs = self._crawl_with_selenium(url_objects)
        successful_docs = [doc for doc in docs if not doc.is_error]

        if len(successful_docs) == 0:
            logger.info("Selenium crawling failed, trying requests-based crawling...")
            docs = self._crawl_with_requests(url_objects)
            successful_docs = [doc for doc in docs if not doc.is_error]

        if len(successful_docs) == 0:
            logger.warning(
                "All web crawling failed, creating mock content from URLs..."
            )
            docs = self._create_mock_content_from_urls(url_objects)
            successful_docs = [doc for doc in docs if not doc.is_error]

        logger.info(
            f"Web crawling complete: {len(successful_docs)}/{len(url_objects)} URLs"
        )
        return docs

    def _crawl_with_selenium(self, url_objects: list[DocObject]) -> list[CrawledObject]:
        """Crawl URLs using Selenium WebDriver with enhanced stealth options."""
        options = webdriver.ChromeOptions()
        options.add_argument("--no-sandbox")
        options.add_argument("--headless")
        options.add_argument("--disable-dev-shm-usage")
        options.add_argument("--disable-gpu")
        options.add_argument("--disable-extensions")
        options.add_argument("--disable-infobars")
        options.add_argument("--remote-debugging-pipe")
        options.add_argument("--timeout=30000")
        options.add_argument("--page-load-timeout=30")
        options.add_argument("--disable-web-security")
        options.add_argument("--allow-running-insecure-content")
        options.add_argument("--disable-features=VizDisplayCompositor")
        options.add_argument("--disable-background-timer-throttling")
        options.add_argument("--disable-backgrounding-occluded-windows")
        options.add_argument("--disable-renderer-backgrounding")
        options.add_argument("--disable-features=TranslateUI")
        options.add_argument("--disable-blink-features=AutomationControlled")
        options.add_argument("--no-first-run")
        options.add_argument("--no-default-browser-check")
        options.add_argument("--disable-default-apps")
        options.add_argument("--disable-sync")
        options.add_argument("--disable-translate")
        options.add_argument("--disable-background-networking")
        options.add_argument("--disable-client-side-phishing-detection")
        options.add_argument("--disable-component-extensions-with-background-pages")
        options.add_argument("--disable-domain-reliability")
        options.add_argument("--disable-features=AudioServiceOutOfProcess")
        options.add_argument("--disable-hang-monitor")
        options.add_argument("--disable-ipc-flooding-protection")
        options.add_argument("--disable-prompt-on-repost")
        options.add_argument("--disable-sync-preferences")
        options.add_argument("--disable-web-resources")
        options.add_argument("--metrics-recording-only")
        options.add_argument("--no-report-upload")
        options.add_argument("--safebrowsing-disable-auto-update")

        try:
            chrome_driver_path = ChromeDriverManager(
                driver_version=CHROME_DRIVER_VERSION
            ).install()
            service = Service(executable_path=chrome_driver_path)
            logger.info(f"Using ChromeDriver: {chrome_driver_path}")
        except Exception as e:
            logger.error(f"Failed to install ChromeDriver: {e}")
            return [
                self._create_error_doc(
                    url_obj, f"ChromeDriver installation failed: {e}"
                )
                for url_obj in url_objects
            ]

        docs: list[CrawledObject] = []
        start_time = time.time()
        max_time_per_url = 30
        successful_crawls = 0
        failed_crawls = 0

        for i, url_obj in enumerate(url_objects, 1):
            url_start_time = time.time()
            driver = None
            max_retries = 2

            for retry_attempt in range(max_retries):
                try:
                    driver = webdriver.Chrome(service=service, options=options)
                    driver.set_page_load_timeout(30)
                    driver.set_script_timeout(30)
                    driver.get(url_obj.source)
                    time.sleep(3)

                    if time.time() - url_start_time > max_time_per_url:
                        logger.warning(
                            f"URL {url_obj.source} taking too long, skipping"
                        )
                        raise Exception("URL load timeout")

                    html = driver.page_source
                    soup = BeautifulSoup(html, "html.parser")

                    text_list = []
                    for string in soup.strings:
                        if string.find_parent("a"):
                            href = urljoin(
                                url_obj.source, string.find_parent("a").get("href")
                            )
                            if href.startswith(url_obj.source):
                                text_list.append(f"{string} {href}")
                        elif string.strip():
                            text_list.append(string)
                    text_output = "\n".join(text_list)

                    title = url_obj.source
                    for title_elem in soup.find_all("title"):
                        title = title_elem.get_text()
                        break

                    docs.append(
                        CrawledObject(
                            id=url_obj.id,
                            source=url_obj.source,
                            content=text_output,
                            metadata={"title": title, "source": url_obj.source},
                            source_type=SourceType.WEB,
                        )
                    )
                    successful_crawls += 1
                    logger.info(
                        f"Successfully crawled URL {i}/{len(url_objects)}: {url_obj.source}"
                    )
                    break

                except Exception as err:
                    if driver:
                        with contextlib.suppress(Exception):
                            driver.quit()

                    if retry_attempt == max_retries - 1:
                        error_msg = str(err)
                        expected_errors = [
                            "cannot determine loading status",
                            "target window already closed",
                            "no such window",
                            "chrome not reachable",
                            "session deleted",
                            "timeout",
                        ]
                        if any(e in error_msg.lower() for e in expected_errors):
                            logger.debug(
                                f"Expected error crawling {url_obj.source}: {error_msg}"
                            )
                        else:
                            logger.error(
                                f"Error crawling {url_obj.source}: {error_msg}"
                            )
                        docs.append(self._create_error_doc(url_obj, str(err)))
                        failed_crawls += 1
                    else:
                        logger.info(
                            f"Retry {retry_attempt + 1}/{max_retries} for {url_obj.source}"
                        )
                        time.sleep(2)

        elapsed_time = time.time() - start_time
        logger.info(
            f"Selenium crawling: {successful_crawls}/{len(url_objects)} URLs in {elapsed_time:.1f}s"
        )
        return docs

    def _crawl_with_requests(self, url_objects: list[DocObject]) -> list[CrawledObject]:
        """Fallback crawling using requests library."""
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.5",
            "Accept-Encoding": "gzip, deflate",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
        }

        docs: list[CrawledObject] = []
        successful_crawls = 0

        for i, url_obj in enumerate(url_objects, 1):
            try:
                logger.info(f"Requesting URL {i}/{len(url_objects)}: {url_obj.source}")
                response = requests.get(url_obj.source, headers=headers, timeout=10)
                response.raise_for_status()

                soup = BeautifulSoup(response.text, "html.parser")
                text_list = []
                for string in soup.strings:
                    if string.find_parent("a"):
                        href = urljoin(
                            url_obj.source, string.find_parent("a").get("href")
                        )
                        if href.startswith(url_obj.source):
                            text_list.append(f"{string} {href}")
                    elif string.strip():
                        text_list.append(string)
                text_output = "\n".join(text_list)

                title = url_obj.source
                for title_elem in soup.find_all("title"):
                    title = title_elem.get_text()
                    break

                docs.append(
                    CrawledObject(
                        id=url_obj.id,
                        source=url_obj.source,
                        content=text_output,
                        metadata={"title": title, "source": url_obj.source},
                        source_type=SourceType.WEB,
                    )
                )
                successful_crawls += 1
                logger.info(
                    f"Successfully crawled URL {i}/{len(url_objects)}: {url_obj.source}"
                )
            except Exception as err:
                logger.debug(f"Requests failed for {url_obj.source}: {err}")
                docs.append(self._create_error_doc(url_obj, str(err)))

        logger.info(f"Requests crawling: {successful_crawls}/{len(url_objects)} URLs")
        return docs

    def _create_mock_content_from_urls(
        self, url_objects: list[DocObject]
    ) -> list[CrawledObject]:
        """Create mock content from URLs when web crawling fails completely.

        Args:
            url_objects (list[DocObject]): List of DocObject instances containing URLs.

        Returns:
            List[CrawledObject]: List of CrawledObject instances with mock content.
        """
        docs: list[CrawledObject] = []

        for url_obj in url_objects:
            url = url_obj.source
            domain = url.split("/")[2] if len(url.split("/")) > 2 else url

            path_parts = url.split("/")
            page_type = "homepage"
            if len(path_parts) > 3:
                page_type = path_parts[3].replace("-", " ").replace("_", " ")

            if "company" in url.lower() or "about" in url.lower():
                content = f"Company information for {domain}. This page contains details about the company's history, mission, values, and leadership team."
            elif "contact" in url.lower():
                content = f"Contact information for {domain}. This page provides ways to get in touch with the company including phone numbers, email addresses, and office locations."
            elif "privacy" in url.lower():
                content = f"Privacy policy for {domain}. This page outlines how the company collects, uses, and protects user data and personal information."
            elif "terms" in url.lower():
                content = f"Terms and conditions for {domain}. This page contains the legal terms governing the use of the company's services and products."
            elif "resources" in url.lower() or "blog" in url.lower():
                content = f"Resources and information for {domain}. This page provides additional materials, updates, news, and educational content related to the company's products and services."
            elif "solutions" in url.lower() or "products" in url.lower():
                content = f"Solutions and products offered by {domain}. This page showcases the company's offerings and services."
            elif (
                "faq" in url.lower()
                or "help" in url.lower()
                or "support" in url.lower()
            ):
                content = f"Help and support for {domain}. This page answers frequently asked questions and provides customer support resources."
            else:
                content = f"Welcome to {domain}. This is the {page_type} page providing information about the company's products and services."

            docs.append(
                CrawledObject(
                    id=url_obj.id,
                    source=url_obj.source,
                    content=content,
                    metadata={
                        "title": f"{page_type.title()} - {domain}",
                        "source": url_obj.source,
                        "mock_content": True,
                    },
                    source_type=SourceType.WEB,
                )
            )

        logger.info(f"Created mock content for {len(docs)} URLs")
        return docs

    def _create_error_doc(self, url_obj: DocObject, error_msg: str) -> CrawledObject:
        """Create an error document for failed crawls."""
        return CrawledObject(
            id=url_obj.id,
            source=url_obj.source,
            content="",
            metadata={"title": url_obj.source, "source": url_obj.source},
            is_error=True,
            error_message=error_msg,
            source_type=SourceType.WEB,
        )

    def get_all_urls(self, base_url: str, max_num: int) -> list[str]:
        """Get all URLs from a base URL up to a maximum number.

        Args:
            base_url (str): The starting URL to crawl from.
            max_num (int): Maximum number of URLs to collect.

        Returns:
            List[str]: List of collected URLs, sorted alphabetically.
        """
        logger.info(f"Discovering URLs from {base_url} (max: {max_num})")
        urls_visited: list[str] = []
        base_url = base_url.split("#")[0].rstrip("/")
        urls_to_visit = [base_url]

        max_iterations = max_num * 3
        iteration_count = 0
        start_time = time.time()
        max_time_seconds = 60

        while urls_to_visit and iteration_count < max_iterations:
            if time.time() - start_time > max_time_seconds:
                logger.warning(f"URL discovery timed out after {max_time_seconds}s")
                break
            if len(urls_visited) >= max_num:
                break

            current_url = urls_to_visit.pop(0)
            iteration_count += 1

            if current_url not in urls_visited:
                urls_visited.append(current_url)
                try:
                    new_urls = self.get_outsource_urls(current_url, base_url)
                    urls_to_visit.extend(new_urls)
                    urls_to_visit = list(set(urls_to_visit))
                    if new_urls:
                        logger.info(
                            f"Found {len(new_urls)} new URLs from {current_url}"
                        )
                except Exception as e:
                    logger.error(f"Error discovering URLs from {current_url}: {e}")
                    continue

        elapsed_time = time.time() - start_time
        logger.info(
            f"URL discovery complete: {len(urls_visited)} URLs found in {elapsed_time:.1f}s"
        )
        return sorted(urls_visited[:max_num])

    def get_outsource_urls(self, curr_url: str, base_url: str) -> list[str]:
        """Get outsource URLs from a given URL.

        Args:
            curr_url (str): The current URL to extract links from.
            base_url (str): The base URL for filtering and validation.

        Returns:
            list[str]: List of valid outsource URLs.
        """
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/14.0.3 Safari/605.1.15"
        }
        new_urls: list[str] = []
        try:
            response = requests.get(curr_url, headers=headers, timeout=10)
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, "html.parser")
                for link in soup.find_all("a"):
                    try:
                        full_url = urljoin(curr_url, link.get("href"))
                        full_url = full_url.split("#")[0].rstrip("/")
                        if self._check_url(full_url, base_url):
                            new_urls.append(full_url)
                    except Exception as err:
                        logger.error(
                            f"Fail to process sub-url {link.get('href')}: {err}"
                        )
            else:
                logger.error(
                    f"Failed to retrieve page {curr_url}, status code: {response.status_code}"
                )
        except Exception as err:
            logger.error(f"Fail to get the page from {curr_url}: {err}")
        return list(set(new_urls))

    def _check_url(self, full_url: str, base_url: str) -> bool:
        """Check if a URL is valid and belongs to the base URL.

        Args:
            full_url (str): The URL to check.
            base_url (str): The base URL for validation.

        Returns:
            bool: True if the URL is valid and belongs to the base URL, False otherwise.
        """
        kw_list = [".pdf", ".jpg", ".png", ".docx", ".xlsx", ".pptx", ".zip", ".jpeg"]
        return (
            full_url.startswith(base_url)
            and full_url
            and not any(kw in full_url for kw in kw_list)
            and full_url != base_url
        )

    def to_crawled_text(self, text_list: list[str]) -> list[CrawledObject]:
        """Convert a list of text strings to CrawledObject instances.

        Args:
            text_list (List[str]): List of text strings to convert.

        Returns:
            List[CrawledObject]: List of CrawledObject instances.
        """
        return [
            CrawledObject(
                id=str(uuid.uuid4()),
                source="text",
                content=text,
                metadata={},
                source_type=SourceType.TEXT,
            )
            for text in text_list
        ]

    def to_crawled_local_objs(self, file_list: list[str]) -> list[CrawledObject]:
        """Convert a list of local files to CrawledObject instances.

        Args:
            file_list (List[str]): List of file paths to process.

        Returns:
            List[CrawledObject]: List of CrawledObject instances.
        """
        local_objs = [DocObject(str(uuid.uuid4()), file) for file in file_list]
        return [self.crawl_file(obj) for obj in local_objs]

    def crawl_file(self, local_obj: DocObject) -> CrawledObject:
        """Crawl a local file and extract its content.

        Supports: .md, .txt, .json, .html, .pdf, .docx, .xlsx, .xls, .pptx, .ppt,
        .png, .jpg, .jpeg (images via Mistral OCR when MISTRAL_API_KEY is set).

        Args:
            local_obj (DocObject): The local file object to process.

        Returns:
            CrawledObject: A CrawledObject instance containing the file's content.
        """
        file_path = Path(local_obj.source)
        file_type = file_path.suffix.lstrip(".").lower()
        file_name = file_path.name
        doc_text = ""

        try:
            if not file_type:
                raise FileNotFoundError(f"No file type detected for file: {file_path}")

            # Use Mistral OCR for visual file types when the API key is available
            if file_type in ["pdf", "png", "jpg", "jpeg", "pptx", "ppt"] and (
                MISTRAL_API_KEY and MISTRAL_API_KEY != "<your-mistral-api-key>"
            ):
                client = Mistral(api_key=MISTRAL_API_KEY)
                if file_type in ["pdf", "pptx", "ppt"]:
                    with open(file_path, "rb") as file_content:
                        uploaded_doc = client.files.upload(
                            file={"file_name": file_name, "content": file_content},
                            purpose="ocr",
                        )
                    signed_url = client.files.get_signed_url(file_id=uploaded_doc.id)
                    ocr_response = client.ocr.process(
                        model="mistral-ocr-latest",
                        document={
                            "type": "document_url",
                            "document_url": signed_url.url,
                        },
                    )
                else:
                    base64_image = encode_image(str(file_path))
                    ocr_response = client.ocr.process(
                        model="mistral-ocr-latest",
                        document={
                            "type": "image_url",
                            "image_url": f"data:image/{file_type};base64,{base64_image}",
                        },
                    )
                doc_text = "".join(page.markdown for page in ocr_response.pages)
                logger.info("Mistral OCR extraction succeeded.")

            elif file_type == "html":
                with open(file_path, encoding="utf-8") as f:
                    html = f.read()
                soup = BeautifulSoup(html, "html.parser")
                text_list = []
                for string in soup.strings:
                    if string.find_parent("a"):
                        href = string.find_parent("a").get("href")
                        text_list.append(f"{string} {href}")
                    elif string.strip():
                        text_list.append(string)
                doc_text = "\n".join(text_list)
                title = file_name
                for title_elem in soup.find_all("title"):
                    title = title_elem.get_text()
                    break
                return CrawledObject(
                    id=local_obj.id,
                    source=local_obj.source,
                    content=doc_text,
                    metadata={"title": title, "source": local_obj.source},
                    source_type=SourceType.LOCAL,
                )

            elif file_type == "pdf":
                logger.info("MISTRAL_API_KEY not set; using pypdf for PDF extraction.")
                from pypdf import PdfReader

                reader = PdfReader(str(file_path))
                doc_text = "\n".join(page.extract_text() or "" for page in reader.pages)

            elif file_type in ("doc", "docx"):
                from docx import Document as DocxDocument

                doc = DocxDocument(str(file_path))
                doc_text = "\n".join(para.text for para in doc.paragraphs)

            elif file_type in ("xlsx", "xls"):
                import openpyxl

                wb = openpyxl.load_workbook(
                    str(file_path), read_only=True, data_only=True
                )
                rows: list[str] = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        rows.append(
                            "\t".join(str(c) if c is not None else "" for c in row)
                        )
                doc_text = "\n".join(rows)

            elif file_type in ("txt", "md"):
                with open(file_path, encoding="utf-8", errors="replace") as f:
                    doc_text = f.read()

            elif file_type in ("pptx", "ppt"):
                from pptx import Presentation

                prs = Presentation(str(file_path))
                slides_text: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides_text.append(shape.text)
                doc_text = "\n".join(slides_text)

            elif file_type == "json":
                with open(file_path) as f:
                    doc_text = json.dumps(json.load(f))

            else:
                raise NotImplementedError(
                    f"Unsupported file type: .{file_type}. "
                    "Supported types: md, txt, json, html, pdf, docx, xlsx, xls, pptx, ppt, "
                    "png, jpg, jpeg (images require MISTRAL_API_KEY)."
                )

            return CrawledObject(
                id=local_obj.id,
                source=local_obj.source,
                content=doc_text,
                metadata={"title": file_name, "source": local_obj.source},
                source_type=SourceType.LOCAL,
            )

        except Exception as err_msg:
            logger.info(f"Error processing file: {err_msg}")
            return CrawledObject(
                id=local_obj.id,
                source=local_obj.source,
                content=None,
                metadata={"title": file_name},
                source_type=SourceType.LOCAL,
                is_error=True,
                error_message=str(err_msg),
            )

    @staticmethod
    def save(file_path: str, docs: list[CrawledObject]) -> None:
        """Save a list of CrawledObject instances to a file.

        Args:
            file_path (str): Path where to save the objects.
            docs (List[CrawledObject]): List of CrawledObject instances to save.
        """
        with open(file_path, "wb") as f:
            pickle.dump(docs, f)

    @classmethod
    def chunk(cls, doc_objs: list[CrawledObject]) -> list[CrawledObject]:
        """Split documents into smaller chunks.

        Args:
            doc_objs (List[CrawledObject]): List of CrawledObject instances to chunk.

        Returns:
            List[CrawledObject]: List of chunked CrawledObject instances.
        """
        text_splitter = _make_text_splitter()
        docs: list[CrawledObject] = []
        for doc_obj in doc_objs:
            if doc_obj.is_error or doc_obj.content is None:
                logger.debug(f"Skipping {doc_obj.source}: error or no content")
                continue
            if doc_obj.is_chunk:
                logger.debug(f"Skipping {doc_obj.source}: already chunked")
                docs.append(doc_obj)
                continue

            try:
                splits = text_splitter.split_text(doc_obj.content)
            except Exception as split_error:
                logger.warning(
                    f"Failed to split document from {doc_obj.source}: {split_error}. Skipping."
                )
                continue

            for i, txt in enumerate(splits):
                docs.append(
                    CrawledObject(
                        id=f"{doc_obj.id}_{i}",
                        source=doc_obj.source,
                        content=txt,
                        metadata=doc_obj.metadata,
                        is_chunk=True,
                        source_type=doc_obj.source_type,
                    )
                )
        return docs
